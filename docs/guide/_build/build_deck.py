"""Build the CryptoCopilot Demo Day deck (.pptx) — a 12-slide, 10-minute talk.

Dark fintech theme to match the app UI; a hand-drawn architecture diagram on
slide 3; screenshots from ../screenshots embedded where they help.
    python3 build_deck.py
"""
import os
from lxml import etree
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

HERE = os.path.dirname(__file__)
SHOTS = os.path.join(HERE, "..", "screenshots")
OUT = os.path.normpath(os.path.join(HERE, "..", "presentation", "CryptoCopilot-DemoDay.pptx"))
SHOT = {
    "signals": "Screenshot 2026-06-06 at 11.39.19.png",
    "markets": "Screenshot 2026-06-06 at 11.39.54.png",
    "ml_pipeline": "Screenshot 2026-06-06 at 11.40.12.png",
    "coin_detail": "Screenshot 2026-06-06 at 11.41.01.png",
    "signal_card": "Screenshot 2026-06-06 at 11.42.10.png",
    "analyst": "Screenshot 2026-06-06 at 11.42.56.png",
    "paper_trades": "Screenshot 2026-06-06 at 11.43.40.png",
    "performance": "Screenshot 2026-06-06 at 11.43.53.png",
    "er_diagram": "Screenshot 2026-06-06 at 12.02.07.png",
}
def shot(a): return os.path.normpath(os.path.join(SHOTS, SHOT[a]))

# ---- palette ----------------------------------------------------------------
BG    = RGBColor(0x0A, 0x10, 0x1C)
CARD  = RGBColor(0x14, 0x21, 0x3A)
CARD2 = RGBColor(0x10, 0x1A, 0x30)
INK   = RGBColor(0xEC, 0xF1, 0xF8)
MUTED = RGBColor(0x9A, 0xA8, 0xBD)
FAINT = RGBColor(0x66, 0x74, 0x8A)
TEAL  = RGBColor(0x2D, 0xD4, 0xBF)
BLUE  = RGBColor(0x5B, 0x9D, 0xF9)
CYAN  = RGBColor(0x38, 0xBD, 0xF8)
ORANGE= RGBColor(0xF2, 0xA6, 0x5A)
GREEN = RGBColor(0x34, 0xD3, 0x99)
RED   = RGBColor(0xF8, 0x71, 0x71)
VIOLET= RGBColor(0xA7, 0x8B, 0xFA)

FONT = "Calibri"
W, H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


# ---- helpers ----------------------------------------------------------------
def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s

def _set_runs(p, runs, align=PP_ALIGN.LEFT, space_after=4, line=1.05):
    p.alignment = align
    p.space_after = Pt(space_after)
    p.line_spacing = line
    for text, size, color, bold in runs:
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = FONT
    return p

def textbox(s, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            space_after=4, line=1.05, wrap=True):
    """lines = list of paragraphs; each paragraph = list of (text,size,color,bold)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2); tf.margin_top = tf.margin_bottom = Pt(2)
    for i, runs in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _set_runs(p, runs, align, space_after, line)
    return tb

def chip(s, x, y, w, h, text, fill, txt=INK, size=11, bold=True, radius=0.10):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background(); sh.shadow.inherit = False
    _round(sh, radius)
    tf = sh.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Pt(6); tf.margin_top = tf.margin_bottom = Pt(2)
    _set_runs(tf.paragraphs[0], [(text, size, txt, bold)], PP_ALIGN.CENTER, 0, 1.0)
    return sh

def card(s, x, y, w, h, accent, title=None, subtitle=None, body=None,
         tsize=15, ssize=10.5, fill=CARD):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = accent; sh.line.width = Pt(1.25); sh.shadow.inherit = False
    _round(sh, 0.08)
    tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = Pt(11); tf.margin_top = Pt(9); tf.margin_bottom = Pt(8)
    first = True
    if title:
        _set_runs(tf.paragraphs[0], [(title, tsize, accent, True)], space_after=3); first = False
    if subtitle:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        _set_runs(p, [(subtitle, ssize, INK, False)], space_after=4, line=1.08); first = False
    for ln in (body or []):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        _set_runs(p, ln, space_after=3, line=1.06); first = False
    return sh

def _round(sh, val):
    try:
        sh.adjustments[0] = val
    except Exception:
        pass

def bar(s, x, y, w, h, color):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background(); sh.shadow.inherit = False
    return sh

def title_block(s, kicker, title, color=TEAL):
    bar(s, 0.6, 0.62, 0.5, 0.12, color)
    textbox(s, 0.62, 0.74, 11.5, 0.4, [[(kicker.upper(), 12, color, True)]])
    textbox(s, 0.6, 1.06, 12.1, 0.95, [[(title, 30, INK, True)]], line=1.0)

def footer(s, idx, note="Decision-support, not financial advice · Paper trading only"):
    bar(s, 0.0, 7.18, 13.333, 0.012, RGBColor(0x24, 0x33, 0x4D))
    textbox(s, 0.6, 7.2, 9.5, 0.3, [[(note, 9, FAINT, False)]])
    textbox(s, 12.2, 7.2, 0.9, 0.3, [[(f"{idx} / 12", 9, FAINT, False)]], align=PP_ALIGN.RIGHT)

def arrow(s, x1, y1, x2, y2, color, width=1.75, dashed=False, both=False):
    cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = color; cn.line.width = Pt(width); cn.shadow.inherit = False
    ln = cn.line._get_or_add_ln()
    if dashed:
        d = etree.SubElement(ln, qn("a:prstDash")); d.set("val", "dash")
    tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", "triangle"); tail.set("w", "med"); tail.set("len", "med")
    if both:
        head = etree.SubElement(ln, qn("a:headEnd"))
        head.set("type", "triangle"); head.set("w", "med"); head.set("len", "med")
    return cn

def picture(s, alias, x, y, max_w, max_h, border=True):
    p = shot(alias); wpx, hpx = Image.open(p).size; ratio = wpx / hpx
    w = max_w; h = w / ratio
    if h > max_h:
        h = max_h; w = h * ratio
    x = x + (max_w - w) / 2; y = y + (max_h - h) / 2
    pic = s.shapes.add_picture(p, Inches(x), Inches(y), Inches(w), Inches(h))
    if border:
        pic.line.color.rgb = RGBColor(0x2A, 0x3A, 0x55); pic.line.width = Pt(1)
    return pic

def bullets(s, x, y, w, h, items, size=14, gap=7, color=INK, marker="▸", mcolor=TEAL):
    lines = []
    for it in items:
        if isinstance(it, tuple):  # (text, color) override
            lines.append([(marker + "  ", size, mcolor, True), (it[0], size, it[1], False)])
        else:
            lines.append([(marker + "  ", size, mcolor, True), (it, size, color, False)])
    return textbox(s, x, y, w, h, lines, space_after=gap, line=1.04)


# ============================================================ SLIDE 1 — TITLE
s = slide()
bar(s, 0.0, 0.0, 13.333, 0.16, TEAL)
chip(s, 0.85, 1.7, 1.5, 1.5, "C", CARD, txt=TEAL, size=64, radius=0.22)
s.shapes[-1].line.color.rgb = TEAL; s.shapes[-1].line.width = Pt(2)
textbox(s, 2.7, 1.75, 10.0, 1.4, [[("CryptoCopilot", 52, INK, True)]], anchor=MSO_ANCHOR.MIDDLE)
textbox(s, 2.72, 3.05, 10.0, 0.6, [[("An intelligent assistant for the world of crypto", 20, TEAL, False)]])
textbox(s, 0.9, 4.3, 11.5, 0.5,
        [[("ML direction signals  ·  technical analysis  ·  cited AI chat  ·  safe paper trading", 16, MUTED, False)]])
textbox(s, 0.9, 4.92, 11.5, 0.5,
        [[("Presented by  ", 16, MUTED, False), ("Kamran Zajkani", 16, TEAL, True)]])
chip(s, 0.9, 5.6, 3.0, 0.5, "Demo Day · June 2026", CARD, txt=INK, size=12)
chip(s, 4.05, 5.6, 4.6, 0.5, "Decision-support, not financial advice", CARD2, txt=MUTED, size=12)
textbox(s, 0.9, 6.5, 11.5, 0.5,
        [[("A polyglot system: Python (data + ML) · Java/Spring Boot · React — five containers, one shared database.",
           12, FAINT, False)]])

# ====================================================== SLIDE 2 — THE HOOK
s = slide()
title_block(s, "Why this project", "Money, Control — and a Different Idea", TEAL)
card(s, 0.6, 2.05, 5.75, 3.7, ORANGE, "Traditional money", None, [
    [("🏦  Controlled by central banks — central control.", 13.5, INK, False)],
    [("💸  You pay fees to hold and move your own money.", 13.5, INK, False)],
    [("🔒  Someone else is always in charge of your account.", 13.5, INK, False)],
], tsize=18)
card(s, 6.6, 2.05, 6.1, 3.7, TEAL, "Cryptocurrency", None, [
    [("🌐  Decentralized — no central authority in control.", 13.5, INK, False)],
    [("🔑  You hold your own money, directly.", 13.5, INK, False)],
    [("➕  Pros: open, global, fewer gatekeepers.", 13.5, GREEN, False)],
    [("➖  Cons: complex, volatile, confusing for newcomers.", 13.5, ORANGE, False)],
], tsize=18)
chip(s, 0.6, 6.05, 12.1, 0.92,
     "CryptoCopilot helps newcomers: it turns scattered market data + news into clear information and signals about coins.",
     CARD, txt=INK, size=15)
footer(s, 2)

# ================================================== SLIDE 3 — ARCHITECTURE
s = slide()
title_block(s, "How it is built", "One System, Two Languages, Five Containers", BLUE)
# central column: frontend -> backend -> db
fx, fw = 4.75, 3.85
card(s, fx, 1.95, fw, 0.82, CYAN, "frontend", "React + nginx — the browser app", tsize=14, ssize=10)
card(s, fx-0.45, 3.12, fw+0.9, 1.18, ORANGE, "backend — Spring Boot",
     "Modular monolith:  ta4j · Spring AI · trading · analyst · REST API", tsize=14, ssize=10)
card(s, fx, 4.95, fw, 1.0, TEAL, "db — Postgres 16 + pgvector",
     "The single shared database", tsize=13.5, ssize=10)
# left: ml
card(s, 0.55, 4.9, 3.55, 1.1, BLUE, "ml — Python",
     "Batch worker: ingestion · XGBoost · SHAP", tsize=13.5, ssize=10)
# right: ml-api
card(s, 9.55, 3.12, 3.45, 1.18, BLUE, "ml-api — Python",
     "FastAPI: on-demand ingest / train / predict", tsize=13.5, ssize=10)
# arrows
arrow(s, fx+fw/2, 2.77, fx+fw/2, 3.12, CYAN)                     # frontend -> backend
textbox(s, fx+fw/2+0.1, 2.78, 3.2, 0.3, [[("REST / JSON (same origin)", 9.5, MUTED, False)]])
arrow(s, fx+fw/2, 4.30, fx+fw/2, 4.95, ORANGE)                   # backend -> db
textbox(s, fx+fw/2+0.1, 4.5, 3.4, 0.3, [[("JDBC — reads ml's tables; writes orders + vectors", 9.5, MUTED, False)]])
arrow(s, 4.10, 5.45, 4.75, 5.45, BLUE)                          # ml -> db (edge to edge)
textbox(s, 0.55, 6.06, 3.6, 0.3, [[("writes predictions + raw data", 9.5, MUTED, False)]])
arrow(s, 9.55, 3.71, 9.05, 3.71, BLUE, both=True)               # backend <-> ml-api
textbox(s, 8.05, 2.78, 2.7, 0.3, [[("HTTP  /api/ml/*", 9.5, MUTED, False)]], align=PP_ALIGN.CENTER)
arrow(s, 10.9, 4.30, 8.55, 5.35, BLUE, dashed=True)             # ml-api -> db
chip(s, 0.55, 6.5, 12.25, 0.62,
     "The database is the polyglot boundary:  Python writes  →  Java reads.   No RPC, no shared model files.",
     CARD, txt=TEAL, size=13.5)
footer(s, 3)

# ================================================== SLIDE 4 — DATA
s = slide()
title_block(s, "The fuel", "Real Data: 10 Coins, 5 Free Sources", GREEN)
coins = "BTC · ETH · SOL · BNB · XRP · ADA · AVAX · DOT · LINK · MATIC"
chip(s, 0.6, 1.95, 12.1, 0.6, coins, CARD, txt=INK, size=16)
card(s, 0.6, 2.78, 6.0, 3.05, BLUE, "Where it comes from", None, [
    [("📈  Binance — price candles (1h / 4h / 1d, ~2 years)", 12.5, INK, False)],
    [("🦎  CoinGecko — market cap, supply, community & dev data", 12.5, INK, False)],
    [("📰  RSS ×5 — CoinDesk, Cointelegraph, Decrypt, The Block…", 12.5, INK, False)],
    [("⛓️  Blockchain.com + Etherscan — on-chain activity", 12.5, INK, False)],
    [("📚  Curated knowledge base — how each coin works", 12.5, INK, False)],
], tsize=16)
card(s, 6.85, 2.78, 5.85, 3.05, TEAL, "How we treat it", None, [
    [("🗃️  ~231,000 rows loaded into one Postgres database", 12.5, INK, False)],
    [("🛡️  Multi-source by design — no single point of failure", 12.5, INK, False)],
    [("🔁  If any source fails: log and skip, never crash", 12.5, INK, False)],
    [("👤  One writer per table — clean, auditable data", 12.5, INK, False)],
    [("💶  All sources are public and free", 12.5, GREEN, False)],
], tsize=16)
footer(s, 4)

# ================================================== SLIDE 5 — ML
s = slide()
title_block(s, "The machine learning", "An Honest Edge, Not a Crystal Ball", VIOLET)
pipe = ["Raw candles", "46 features", "XGBoost", "Calibration", "SHAP", "Probabilities"]
pcols = [BLUE, BLUE, VIOLET, TEAL, ORANGE, GREEN]
px, pw, gap = 0.6, 1.83, 0.18
for i, (t, c) in enumerate(zip(pipe, pcols)):
    chip(s, px + i*(pw+gap), 1.95, pw, 0.62, t, CARD, txt=c, size=12.5)
    if i < len(pipe)-1:
        arrow(s, px + i*(pw+gap)+pw, 2.26, px + (i+1)*(pw+gap), 2.26, FAINT, 1.4)
bullets(s, 0.62, 2.95, 6.6, 3.4, [
    "Predicts UP / FLAT / DOWN over the next 24 hours (a ±2% move).",
    ("ROC-AUC = 0.578 — a small but real edge above luck (0.5 = a coin flip).", INK),
    "Calibrated probabilities: “70%” really means about 70%.",
    "SHAP explains the top-3 reasons behind every single prediction.",
    "Trained on ~2 years of history — the score is reported honestly.",
], size=14, gap=11)
picture(s, "signal_card", 7.5, 2.9, 5.2, 3.6)
footer(s, 5)

# ================================================== SLIDE 6 — BACKEND / TA
s = slide()
title_block(s, "Technical analysis", "Chart Signals You Can Audit", ORANGE)
bullets(s, 0.62, 2.0, 6.3, 3.6, [
    "Computed in Java with ta4j, straight from the price candles.",
    "Indicators: Ichimoku · RSI · MACD · Bollinger · ATR.",
    "Deterministic rules → BULLISH / NEUTRAL / BEARISH + a confidence.",
    "It shows the exact rules that fired — fully transparent.",
    "Independent of the ML model — a real second opinion.",
], size=14, gap=11)
chip(s, 0.62, 5.7, 6.3, 0.7, "Backend = one modular monolith (Spring Boot), not microservices.",
     CARD, txt=ORANGE, size=12.5)
picture(s, "signals", 7.2, 1.95, 5.5, 4.6)
footer(s, 6)

# ================================================== SLIDE 7 — RAG
s = slide()
title_block(s, "The Researcher", "AI Chat That Cites Its Sources", CYAN)
bullets(s, 0.62, 2.0, 6.5, 3.7, [
    "Embeddings turn text into numbers that capture meaning.",
    "pgvector stores them inside the same Postgres database.",
    "RAG = retrieve real sources, then answer using only them.",
    ("Every answer is cited [1][2]; it refuses if the answer isn’t there.", INK),
    "Never gives trading advice — by design.",
    ("recall@8 = 0.90 · 100% of answers cited.", GREEN),
], size=13.5, gap=9)
chip(s, 0.62, 6.0, 6.5, 0.78,
     "Runs free & local (Ollama llama3.2:3b) — or switch to OpenAI gpt-4o-mini with one click.",
     CARD, txt=CYAN, size=12)
picture(s, "analyst", 7.35, 2.0, 5.4, 4.5)
footer(s, 7)

# ================================================== SLIDE 8 — ANALYST
s = slide()
title_block(s, "Putting it together", "Four Views → One Clear Opinion", TEAL)
quad = [("🤖  ML signal", VIOLET), ("📊  Technical", ORANGE), ("🏥  Fundamentals", BLUE), ("📰  News", CYAN)]
for i, (t, c) in enumerate(quad):
    chip(s, 0.6 + i*1.62, 1.98, 1.5, 0.62, t, CARD, txt=c, size=12)
arrow(s, 3.4, 2.62, 3.4, 3.0, TEAL, 2)
chip(s, 0.6, 3.05, 6.25, 0.7, "Analyst  →  Direction · Conviction · Agreement", CARD2, txt=TEAL, size=14)
bullets(s, 0.62, 4.0, 6.3, 2.7, [
    "A simple, deterministic score from −2…+2 per view.",
    "Same inputs always give the same verdict — auditable.",
    "A hallucination guard keeps the summary truthful.",
    "Works even with the AI switched off (plain template).",
], size=13, gap=8)
picture(s, "analyst", 7.2, 1.95, 5.5, 4.6)
footer(s, 8)

# ================================================== SLIDE 9 — PAPER TRADING
s = slide()
title_block(s, "Acting on it — safely", "Practice Trading, No Real Money", GREEN)
bullets(s, 0.62, 2.0, 5.0, 3.6, [
    "Paper trading = real mechanics, fake money.",
    "Long-only; realistic fills with slippage + a 0.1% fee.",
    "Market & limit orders; positions, trades, account tracked.",
    "Performance: equity curve + Sharpe, Sortino, drawdown.",
    ("Honest backtest result — the engine is correct, not magic.", INK),
], size=13.5, gap=9)
picture(s, "paper_trades", 5.75, 1.95, 3.55, 4.3)
picture(s, "performance", 9.45, 2.55, 3.3, 3.1)
footer(s, 9)

# ================================================== SLIDE 10 — THE APP
s = slide()
title_block(s, "The product", "A Clean, Typed Web App", CYAN)
bullets(s, 0.62, 2.0, 5.0, 3.3, [
    "React · Vite · TypeScript · nginx.",
    "Candlestick charts with klinecharts; equity with Recharts.",
    "7 pages: Markets, Signals, Analyst, Chat, Trades, Performance…",
    "…plus an ML Pipeline page — run the whole pipeline in-browser.",
    "A disclaimer is on every page.",
], size=13.5, gap=9)
picture(s, "markets", 5.7, 1.95, 7.0, 2.25)
picture(s, "ml_pipeline", 5.7, 4.35, 7.0, 2.25)
footer(s, 10)

# ================================================== SLIDE 11 — RESULTS
s = slide()
title_block(s, "The honest scorecard", "What It Achieves — Reported Honestly", VIOLET)
card(s, 0.6, 2.0, 3.95, 2.5, VIOLET, "Machine learning", None, [
    [("ROC-AUC  0.578", 18, INK, True)],
    [("A small, real edge above luck.", 11.5, MUTED, False)],
    [("AUC + Brier pass; F1 limited by ~2y data.", 11, MUTED, False)],
], tsize=15)
card(s, 4.7, 2.0, 3.95, 2.5, CYAN, "RAG retrieval", None, [
    [("recall@8  0.90", 18, INK, True)],
    [("100% of answers cited.", 11.5, MUTED, False)],
    [("Refuses the unknown; ≈ €0, local.", 11, MUTED, False)],
], tsize=15)
card(s, 8.8, 2.0, 3.95, 2.5, GREEN, "Engineering", None, [
    [("5 containers · 1 DB", 18, INK, True)],
    [("Tests + CI across 3 languages.", 11.5, MUTED, False)],
    [("Backtest honest, engine correct.", 11, MUTED, False)],
], tsize=15)
chip(s, 0.6, 4.85, 12.15, 0.78,
     "The goal was production-grade polyglot engineering — built end-to-end and reported honestly.",
     CARD, txt=INK, size=15)
bullets(s, 0.62, 5.85, 12.0, 1.1, [
    "Hard rules, never broken: paper only · crypto only · no shorts/leverage · decision-support, not advice.",
], size=12.5, gap=4, mcolor=ORANGE)
footer(s, 11)

# ================================================== SLIDE 12 — THANK YOU
s = slide()
bar(s, 0.0, 0.0, 13.333, 0.16, TEAL)
textbox(s, 0.9, 2.3, 11.5, 1.0, [[("Thank you!", 46, INK, True)]], anchor=MSO_ANCHOR.MIDDLE)
textbox(s, 0.92, 3.4, 11.5, 0.6, [[("Questions?", 24, TEAL, False)]])
textbox(s, 0.9, 4.35, 11.5, 0.6,
        [[("CryptoCopilot — four views on every coin, fused into one explainable opinion.", 16, MUTED, False)]])
chip(s, 0.9, 5.3, 3.7, 0.5, "app  ·  localhost:3000", CARD, txt=INK, size=12)
chip(s, 4.75, 5.3, 4.1, 0.5, "API docs  ·  :8080/swagger-ui.html", CARD, txt=INK, size=12)
chip(s, 9.0, 5.3, 3.4, 0.5, "ML API  ·  :8000/docs", CARD, txt=INK, size=12)
textbox(s, 0.9, 6.5, 11.5, 0.5,
        [[("Decision-support, not financial advice. Paper trading only — no real money, ever.", 12, FAINT, False)]])

prs.save(OUT)
print(f"wrote {OUT}  ({len(prs.slides)} slides)")
